"""
Document Content Parser

Extracts text content from various document formats:
- PDF (.pdf)
- Word (.docx, .doc)
- Excel (.xlsx, .xls)
- PowerPoint (.pptx)
- OpenDocument (.odt, .ods, .odp)
- RTF (.rtf)
- EPUB (.epub)
- Email (.eml, .msg)
- HTML (.html, .htm) - with tag stripping
- Images with OCR (.png, .jpg, .jpeg, .tiff, .bmp, .gif, .webp)
- Plain text (.txt, .md, .rst, .json, .xml, .csv)
- Code files (.py, .js, .ts, etc.)
"""

import io
from pathlib import Path
from typing import Optional

from ragtime.core.file_constants import DOCUMENT_EXTENSIONS, OCR_EXTENSIONS
from ragtime.core.logging import get_logger

logger = get_logger(__name__)


def extract_text_from_file(
    file_path: Path,
    content: Optional[bytes] = None,
    enable_ocr: bool = False,
) -> str:
    """
    Extract text content from a file based on its extension.

    Args:
        file_path: Path to the file
        content: Optional pre-loaded file content (bytes)
        enable_ocr: Whether to use OCR for image files (slower but extracts text from images)

    Returns:
        Extracted text content as string
    """
    suffix = file_path.suffix.lower()

    # Load content if not provided
    if content is None:
        try:
            content = file_path.read_bytes()
        except Exception as e:
            logger.warning(f"Failed to read file {file_path}: {e}")
            return ""

    # Route to appropriate parser
    try:
        if suffix == ".pdf":
            return _extract_pdf(content)
        elif suffix == ".docx":
            return _extract_docx(content)
        elif suffix == ".doc":
            return _extract_doc_legacy(file_path, content)
        elif suffix == ".xlsx":
            return _extract_xlsx(content)
        elif suffix == ".xls":
            return _extract_xls(content)
        elif suffix == ".pptx":
            return _extract_pptx(content)
        elif suffix == ".odt":
            return _extract_odt(content)
        elif suffix == ".ods":
            return _extract_ods(content)
        elif suffix == ".odp":
            return _extract_odp(content)
        elif suffix == ".rtf":
            return _extract_rtf(content)
        elif suffix == ".epub":
            return _extract_epub(content)
        elif suffix == ".eml":
            return _extract_eml(content)
        elif suffix == ".msg":
            return _extract_msg(content)
        elif suffix in {".html", ".htm"}:
            return _extract_html(content)
        elif suffix in OCR_EXTENSIONS and enable_ocr:
            return _extract_image_ocr(content)
        elif suffix in OCR_EXTENSIONS:
            # OCR disabled, skip image files
            logger.debug(f"Skipping image {file_path.name} - OCR disabled")
            return ""
        else:
            # Plain text files
            return _extract_text(content)
    except Exception as e:
        logger.warning(f"Failed to extract text from {file_path}: {e}")
        return ""


def _extract_pdf(content: bytes) -> str:
    """Extract text from PDF file."""
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.warning("pypdf not installed, cannot extract PDF content")
        return ""

    try:
        reader = PdfReader(io.BytesIO(content))
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"PDF extraction error: {e}")
        return ""


def _extract_docx(content: bytes) -> str:
    """Extract text from Word DOCX file."""
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed, cannot extract DOCX content")
        return ""

    try:
        doc = Document(io.BytesIO(content))
        text_parts = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    text_parts.append(row_text)

        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"DOCX extraction error: {e}")
        return ""


def _extract_doc_legacy(file_path: Path, content: bytes) -> str:
    """Extract text from legacy Word DOC file."""
    # Legacy .doc files are more complex
    # Try antiword if available, otherwise skip
    import subprocess

    try:
        result = subprocess.run(
            ["antiword", str(file_path)], capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return result.stdout
    except FileNotFoundError:
        logger.debug("antiword not installed, cannot extract legacy DOC content")
    except Exception as e:
        logger.warning(f"DOC extraction error: {e}")

    return ""


def _extract_xlsx(content: bytes) -> str:
    """Extract text from Excel XLSX file."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed, cannot extract XLSX content")
        return ""

    try:
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"## Sheet: {sheet_name}")

            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        row_values.append(str(cell.value))
                if row_values:
                    text_parts.append(" | ".join(row_values))

        wb.close()
        return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"XLSX extraction error: {e}")
        return ""


def _extract_xls(content: bytes) -> str:
    """Extract text from legacy Excel XLS file."""
    try:
        import xlrd
    except ImportError:
        logger.warning("xlrd not installed, cannot extract XLS content")
        return ""

    try:
        wb = xlrd.open_workbook(file_contents=content)
        text_parts = []

        for sheet in wb.sheets():
            text_parts.append(f"## Sheet: {sheet.name}")

            for row_idx in range(sheet.nrows):
                row_values = []
                for col_idx in range(sheet.ncols):
                    cell = sheet.cell(row_idx, col_idx)
                    if cell.value:
                        row_values.append(str(cell.value))
                if row_values:
                    text_parts.append(" | ".join(row_values))

        return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"XLS extraction error: {e}")
        return ""


def _extract_pptx(content: bytes) -> str:
    """Extract text from PowerPoint PPTX file."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed, cannot extract PPTX content")
        return ""

    try:
        prs = Presentation(io.BytesIO(content))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)

            if slide_texts:
                text_parts.append(f"## Slide {slide_num}")
                text_parts.extend(slide_texts)

        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"PPTX extraction error: {e}")
        return ""


def _extract_odt(content: bytes) -> str:
    """Extract text from OpenDocument Text file."""
    try:
        from odf import text as odf_text
        from odf.opendocument import load
    except ImportError:
        logger.warning("odfpy not installed, cannot extract ODT content")
        return ""

    try:
        doc = load(io.BytesIO(content))
        text_parts = []

        for para in doc.getElementsByType(odf_text.P):
            text = ""
            for node in para.childNodes:
                if hasattr(node, "data"):
                    text += node.data
            if text.strip():
                text_parts.append(text)

        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"ODT extraction error: {e}")
        return ""


def _extract_ods(content: bytes) -> str:
    """Extract text from OpenDocument Spreadsheet file."""
    try:
        from odf import table as odf_table
        from odf import text as odf_text
        from odf.opendocument import load
    except ImportError:
        logger.warning("odfpy not installed, cannot extract ODS content")
        return ""

    try:
        doc = load(io.BytesIO(content))
        text_parts = []

        for sheet in doc.getElementsByType(odf_table.Table):
            sheet_name = sheet.getAttribute("name") or "Sheet"
            text_parts.append(f"## Sheet: {sheet_name}")

            for row in sheet.getElementsByType(odf_table.TableRow):
                row_values = []
                for cell in row.getElementsByType(odf_table.TableCell):
                    cell_text = ""
                    for p in cell.getElementsByType(odf_text.P):
                        for node in p.childNodes:
                            if hasattr(node, "data"):
                                cell_text += node.data
                    if cell_text:
                        row_values.append(cell_text)
                if row_values:
                    text_parts.append(" | ".join(row_values))

        return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"ODS extraction error: {e}")
        return ""


def _extract_odp(content: bytes) -> str:
    """Extract text from OpenDocument Presentation file."""
    try:
        from odf import draw as odf_draw
        from odf import text as odf_text
        from odf.opendocument import load
    except ImportError:
        logger.warning("odfpy not installed, cannot extract ODP content")
        return ""

    try:
        doc = load(io.BytesIO(content))
        text_parts = []

        for page_num, page in enumerate(doc.getElementsByType(odf_draw.Page), 1):
            page_texts = []
            for frame in page.getElementsByType(odf_draw.Frame):
                for text_box in frame.getElementsByType(odf_draw.TextBox):
                    for p in text_box.getElementsByType(odf_text.P):
                        text = ""
                        for node in p.childNodes:
                            if hasattr(node, "data"):
                                text += node.data
                        if text.strip():
                            page_texts.append(text)

            if page_texts:
                text_parts.append(f"## Slide {page_num}")
                text_parts.extend(page_texts)

        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"ODP extraction error: {e}")
        return ""


def _extract_text(content: bytes) -> str:
    """Extract text from plain text file."""
    # Try common encodings
    for encoding in ["utf-8", "latin-1", "cp1252", "ascii"]:
        try:
            return content.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue

    # Last resort: decode with errors replaced
    return content.decode("utf-8", errors="replace")


def _extract_rtf(content: bytes) -> str:
    """Extract text from RTF file."""
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        logger.warning("striprtf not installed, cannot extract RTF content")
        return ""

    try:
        # Decode RTF content
        text = content.decode("utf-8", errors="replace")
        return rtf_to_text(text)
    except Exception as e:
        logger.warning(f"RTF extraction error: {e}")
        return ""


def _extract_epub(content: bytes) -> str:
    """Extract text from EPUB ebook file."""
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError:
        logger.warning("ebooklib not installed, cannot extract EPUB content")
        return ""

    try:
        from bs4 import BeautifulSoup
    except ImportError:
        logger.warning("beautifulsoup4 not installed, cannot extract EPUB content")
        return ""

    try:
        book = epub.read_epub(io.BytesIO(content))
        text_parts = []

        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text = soup.get_text(separator="\n", strip=True)
                if text:
                    text_parts.append(text)

        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"EPUB extraction error: {e}")
        return ""


def _extract_eml(content: bytes) -> str:
    """Extract text from email EML file."""
    import email
    from email.policy import default

    try:
        msg = email.message_from_bytes(content, policy=default)
        text_parts = []

        # Add headers
        if msg.get("Subject"):
            text_parts.append(f"Subject: {msg.get('Subject')}")
        if msg.get("From"):
            text_parts.append(f"From: {msg.get('From')}")
        if msg.get("To"):
            text_parts.append(f"To: {msg.get('To')}")
        if msg.get("Date"):
            text_parts.append(f"Date: {msg.get('Date')}")

        text_parts.append("")  # Blank line after headers

        # Extract body
        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                if content_type == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or "utf-8"
                        text_parts.append(payload.decode(charset, errors="replace"))
                elif content_type == "text/html":
                    payload = part.get_payload(decode=True)
                    if payload:
                        # Strip HTML tags for plain text
                        try:
                            from bs4 import BeautifulSoup

                            charset = part.get_content_charset() or "utf-8"
                            soup = BeautifulSoup(
                                payload.decode(charset, errors="replace"), "html.parser"
                            )
                            text_parts.append(soup.get_text(separator="\n", strip=True))
                        except ImportError:
                            # Fallback: just decode and append
                            charset = part.get_content_charset() or "utf-8"
                            text_parts.append(payload.decode(charset, errors="replace"))
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                charset = msg.get_content_charset() or "utf-8"
                text_parts.append(payload.decode(charset, errors="replace"))

        return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"EML extraction error: {e}")
        return ""


def _extract_msg(content: bytes) -> str:
    """Extract text from Outlook MSG file."""
    try:
        import extract_msg
    except ImportError:
        logger.warning("extract-msg not installed, cannot extract MSG content")
        return ""

    try:
        msg = extract_msg.openMsg(io.BytesIO(content))
        text_parts = []

        # Add headers
        if msg.subject:
            text_parts.append(f"Subject: {msg.subject}")
        if msg.sender:
            text_parts.append(f"From: {msg.sender}")
        if msg.to:
            text_parts.append(f"To: {msg.to}")
        if msg.date:
            text_parts.append(f"Date: {msg.date}")

        text_parts.append("")  # Blank line after headers

        # Add body
        if msg.body:
            text_parts.append(msg.body)

        msg.close()
        return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"MSG extraction error: {e}")
        return ""


def _extract_html(content: bytes) -> str:
    """Extract text from HTML file, stripping tags."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        # Fallback to plain text extraction
        return _extract_text(content)

    try:
        # Decode content
        text = _extract_text(content)

        # Parse and extract text
        soup = BeautifulSoup(text, "html.parser")

        # Remove script and style elements
        for script in soup(["script", "style", "meta", "link"]):
            script.decompose()

        # Get text with newlines preserved
        return soup.get_text(separator="\n", strip=True)
    except Exception as e:
        logger.warning(f"HTML extraction error: {e}")
        return _extract_text(content)


def _extract_image_ocr(content: bytes) -> str:
    """Extract text from image using OCR (Tesseract)."""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        logger.warning(
            "pytesseract/Pillow not installed, cannot perform OCR. "
            "Install with: pip install pytesseract Pillow"
        )
        return ""

    try:
        import time

        start_time = time.time()

        # Open image from bytes
        image = Image.open(io.BytesIO(content))
        original_size = image.size

        # Aggressively resize for fast OCR - 1200px is plenty for text extraction
        max_dimension = 1200
        if max(image.size) > max_dimension:
            ratio = max_dimension / max(image.size)
            new_size = tuple(int(dim * ratio) for dim in image.size)
            image = image.resize(new_size, Image.Resampling.LANCZOS)
            logger.debug(f"OCR: Resized image from {original_size} to {image.size}")
        else:
            logger.debug(f"OCR: Image size {image.size} within limit, no resize needed")

        # Perform OCR with 10s timeout
        import subprocess

        text = pytesseract.image_to_string(image, timeout=10)
        elapsed = time.time() - start_time
        logger.debug(f"OCR completed in {elapsed:.2f}s, extracted {len(text)} chars")

        return text.strip()
    except subprocess.TimeoutExpired:
        logger.warning("OCR timeout (10s) - skipping image")
        return ""
    except Exception as e:
        # Check if Tesseract is not installed
        if "tesseract is not installed" in str(e).lower():
            logger.warning(
                "Tesseract OCR not installed on system. "
                "Install with: apt-get install tesseract-ocr"
            )
        else:
            logger.warning(f"OCR extraction error: {e}")
        return ""


def is_supported_document(file_path: Path) -> bool:
    """Check if a file type is supported for text extraction."""
    return file_path.suffix.lower() in DOCUMENT_EXTENSIONS


def is_ocr_supported(file_path: Path) -> bool:
    """Check if a file type supports OCR extraction."""
    return file_path.suffix.lower() in OCR_EXTENSIONS
